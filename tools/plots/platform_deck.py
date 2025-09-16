from __future__ import annotations
from pathlib import Path
from typing import Any
from pptx import Presentation as _Presentation
from pptx.util import Inches


def _add_title_slide(prs: Any, title: str, subtitle: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle


def _add_image_slide(
    prs: Any, title: str, image_path: Path, height_in: float = 6.0
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.2)
    height = Inches(height_in)
    slide.shapes.add_picture(str(image_path), left, top, height=height)


def build_platform_deck(
    out_pptx: Path | str = "artifacts/pitch/platform/platform_overview.pptx",
) -> Path:
    base = Path("artifacts/pitch/platform")
    prs = _Presentation()
    _add_title_slide(
        prs,
        "General-Purpose In Vitro–Integrated AI Platform",
        "Disease-agnostic core with thin disease adapters and spinout pathways",
    )

    # Conceptual KG
    kg_img = base / "conceptual_platform_kg.png"
    if kg_img.exists():
        _add_image_slide(prs, "Conceptual Platform Knowledge Graph", kg_img)

    # Architecture overview
    arch_img = base / "platform_architecture_overview.png"
    if arch_img.exists():
        _add_image_slide(prs, "Platform Architecture with Disease Adapters", arch_img)

    # Capabilities matrix
    cap_img = base / "capabilities_matrix.png"
    if cap_img.exists():
        _add_image_slide(prs, "Capabilities × Product Opportunities", cap_img)

    # Spinout pathways
    spin_img = base / "spinout_pathways.png"
    if spin_img.exists():
        _add_image_slide(prs, "Spinout Pathways from the Core Platform", spin_img)

    out = Path(out_pptx)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"[deck] Wrote platform deck to {out}")
    return out


if __name__ == "__main__":
    build_platform_deck()
