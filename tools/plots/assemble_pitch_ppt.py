#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Assemble a PPTX pitch deck from generated figures in artifacts/pitch/.        )

    # 7) Roadmap  2) Conceptual KG (image + caption)
  3) Competitive positioning (image + caption)
  4) Demo highlight (image + caption)
  5) Pipeline overview (image, if available)
  6) KG at-a-glance (image, if available)
  7) Roadmap (bullets)
  8) Closing
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from typing import Any

PITCH = Path("artifacts/pitch")
OUT = PITCH / "pitch_deck.pptx"


def add_title(prs: Any, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle


def add_image_slide(prs: Any, title: str, image_path: Path, caption: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = title
    if image_path.exists():
        left = Inches(0.7)
        top = Inches(1.3)
        width = Inches(8.6)
        slide.shapes.add_picture(str(image_path), left, top, width=width)
        if caption:
            txbox = slide.shapes.add_textbox(
                Inches(0.7), Inches(6.8), Inches(8.6), Inches(0.7)
            )
            p = txbox.text_frame.paragraphs[0]
            p.text = caption
            p.font.size = Pt(12)
    else:
        body = slide.shapes.add_textbox(
            Inches(0.7), Inches(2.0), Inches(8.6), Inches(1.0)
        )
        p = body.text_frame.paragraphs[0]
        p.text = f"Missing image: {image_path}"
        p.font.size = Pt(14)


def add_bullets_slide(prs: Any, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.level = 0
        p.font.size = Pt(18)


def main():
    PITCH.mkdir(parents=True, exist_ok=True)
    prs = Presentation()

    # 1) Title
    add_title(
        prs, "Causal + Associative Knowledge Graph Platform", "AKI demo and roadmap"
    )

    # 2) Enhanced Conceptual KG
    add_image_slide(
        prs,
        title="Enhanced conceptual knowledge graph",
        image_path=PITCH / "enhanced_conceptual_kg.png",
        caption=(
            "Comprehensive knowledge graph integrating 5 data sources (OpenTargets, GTEx, HPA, ENCODE, PRIDE) "
            "with cross-layer connections and mechanistic pathways for AKI biomarker discovery."
        ),
    )

    # 3) Industry AI Pipeline Comparison
    add_image_slide(
        prs,
        title="Industry AI pipeline comparison",
        image_path=PITCH / "experimental_rigor_comparison.png",
        caption=(
            "Competitive analysis vs. industry leaders (Recursion, Insilico, BenevolentAI, DeepMind, Atomwise). "
            "Shows balanced approach across experimental integration, mechanistic understanding, and clinical translation."
        ),
    )

    # 4) Realistic Pipeline Overview
    add_image_slide(
        prs,
        title="Realistic pipeline overview",
        image_path=PITCH / "realistic_pipeline_overview.png",
        caption=(
            "End-to-end workflow from multi-source data integration through knowledge graph construction "
            "to biomarker ranking and validation planning with quantified metrics."
        ),
    )

    # 5) Pipeline overview (optional legacy)
    if (PITCH / "pipeline_overview.png").exists():
        add_image_slide(
            prs,
            title="Pipeline overview (alternative)",
            image_path=PITCH / "pipeline_overview.png",
            caption="High-level flow: data → KG → ranking → experiment planning → reports/cards.",
        )

    # 6) Performance metrics
    if (PITCH / "precision_at_k.png").exists():
        add_image_slide(
            prs,
            title="Benchmarking results",
            image_path=PITCH / "precision_at_k.png",
            caption="Precision@K evaluation showing pipeline performance on AKI biomarker discovery.",
        )

    # 7) Roadmap
    add_bullets_slide(
        prs,
        title="Roadmap (next 30–60 days)",
        bullets=[
            "Causal uplift: DML, MR; signed causal paths",
            "Experiment planner: identifiability + VoI/EIG",
            "Ontology alignment: HGNC/EFO/HPO; CURIEs; Biolink mapping",
            "Personalization: cohorts and timepoint/stage propagation",
            "Scale + ops: graph store option, UI, CI/DVC",
        ],
    )

    # 8) Closing
    add_bullets_slide(
        prs,
        title="Get in touch",
        bullets=[
            "Pilot a new indication (we’ll wire new data sources)",
            "Co-design experiments using the planner",
            "Extend the KG with your private data",
        ],
    )

    prs.save(str(OUT))
    print(f"[pitch:ppt] Wrote {OUT}")


if __name__ == "__main__":
    main()
